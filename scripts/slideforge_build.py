# -*- coding: utf-8 -*-
"""
slideforge_build.py — core engine for the slideforge skill.

Takes a structured "scene" JSON and produces a NATIVE editable .pptx. Every object
is a real PowerPoint object (text box, autoshape, table, connector, native chart,
gradient fill), never a rasterized page. Reusable across RECREATE and DESIGN modes.

Scene shape (see references/recreate-contract.md and examples/):

{
  "canvas": "16:9",                 // or "9:16" / "3:4"
  "slides": [
    {
      "background": {"type":"gradient","color1":"#FBFCFF","color2":"#EDF3FA","angle":90}
                    | {"type":"solid","color":"#FFFFFF"},
      "elements": [
        {"type":"text",     "x":0.6,"y":0.4,"w":12,"h":0.5,"size":30,"color":"#15375F",
         "bold":true,"align":"left","content":"Title"},
        {"type":"box",      "x":0.6,"y":1.5,"w":4,"h":2.5,"fill":"#F0F5FB","line":"#AEB8C6",
         "radius":0.08,"content":[["Line 1",11,"#172033",false]]},
        {"type":"table",    "x":4,"y":2.6,"w":3.4,"h":1.8,
         "header":["A","B"],"rows":[["1","x"],["2","y"]],
         "rowfills":[["#FFF0E2","#FFF0E2"],["#E4F4E9","#E4F4E9"]]},
        {"type":"connector","x1":4,"y1":2.2,"x2":6,"y2":2.2,"color":"#2F6FB2","w":2.5,"arrow":true},
        {"type":"bar_chart","x":8,"y":2,"w":4,"h":3,"categories":["A","B","C"],
         "series":[{"name":"Sales","values":[1,2,3]}],"color":"#2F6FB2"},
        {"type":"icon",     "x":1,"y":2,"w":0.5,"h":0.5,"name":"database","color":"#2F6FB2"},
        {"type":"image",    "x":3,"y":2,"w":4,"h":2.5,"path":"...png"},
        {"type":"note",     "x":0.6,"y":6.9,"w":12,"h":0.4,"size":10,"color":"#5F6B7A",
         "content":"Data source: ..."}
      ]
    }
  ]
}

Run:  python slideforge_build.py scene.json -o out.pptx [--canvas 16:9]
"""
from __future__ import annotations

import argparse, json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn

FONT = "Calibri"

CANVAS = {
    "16:9": (13.333, 7.5),
    "9:16": (7.5, 13.333),
    "3:4":  (7.5, 10.0),
    "4:3":  (10.0, 7.5),
}


def _rgb(hexstr: str | None, default="000000") -> RGBColor:
    if not hexstr:
        hexstr = default
    h = hexstr.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _in(v):
    return Inches(v)


def _apply_fill(shape, spec):
    if spec is None:
        shape.fill.background()
        return
    if isinstance(spec, dict):
        if spec.get("type") == "gradient":
            shape.fill.gradient()
            st = shape.fill.gradient_stops
            st[0].color.rgb = _rgb(spec.get("color1", "#FFFFFF"))
            st[0].position = 0.0
            st[1].color.rgb = _rgb(spec.get("color2", spec.get("color1", "#FFFFFF")))
            st[1].position = 1.0
            try:
                shape.fill.gradient_angle = float(spec.get("angle", 90))
            except Exception:
                pass
            return
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(spec.get("color", "#FFFFFF"))
        return
    if spec in ("none",):
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(spec, "#FFFFFF")


def _apply_line(shape, spec, w=1.0):
    if spec is None or spec == "none":
        shape.line.fill.background()
        return
    shape.line.color.rgb = _rgb(spec, "#AEB8C6")
    shape.line.width = Pt(w)


def _set_text(tf, content, size=12, color="#172033", bold=False, align="left",
              anchor="top", wrap=True):
    tf.word_wrap = wrap
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}.get(anchor, MSO_ANCHOR.TOP)
    al = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    if isinstance(content, str):
        # keep raw hex string so the run loop converts once
        paras = [[(content, size, color, bold)]]
    else:
        paras = content  # list of paragraphs; each is list of (text,size,hexcolor,bold)
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = al
        p.space_after = Pt(3)
        # para is either (a) a single run tuple (text,size,color,bold)
        #           or (b) a list of run tuples
        if isinstance(para, (list, tuple)) and len(para) == 4 and isinstance(para[0], str):
            runs = [para]
        else:
            runs = para
        for item in runs:
            text, sz, col, b = item
            r = p.add_run()
            r.text = text
            r.font.size = Pt(sz)
            r.font.color.rgb = _rgb(col, "#172033")
            r.font.bold = b
            r.font.name = FONT


def _add_textbox(slide, e):
    tb = slide.shapes.add_textbox(_in(e["x"]), _in(e["y"]), _in(e["w"]), _in(e["h"]))
    _set_text(tb.text_frame, e.get("content", ""), e.get("size", 12), e.get("color", "#172033"),
              e.get("bold", False), e.get("align", "left"), e.get("anchor", "top"))
    return tb


def _add_box(slide, e):
    shapename = e.get("shape")
    if shapename == "oval":
        shape = MSO_SHAPE.OVAL
    elif shapename == "diamond":
        shape = MSO_SHAPE.DIAMOND
    elif shapename == "round":
        shape = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        shape = MSO_SHAPE.ROUNDED_RECTANGLE if e.get("radius") is not None else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape, _in(e["x"]), _in(e["y"]), _in(e["w"]), _in(e["h"]))
    sp.shadow.inherit = False
    _apply_fill(sp, e.get("fill", "#FFFFFF"))
    _apply_line(sp, e.get("line"), e.get("line_w", 1.0))
    if e.get("radius") is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = e["radius"]
        except Exception:
            pass
    content = e.get("content")
    if content:
        _set_text(sp.text_frame, content, e.get("font_size", 11),
                  e.get("font_color", "#172033"), e.get("bold", False),
                  e.get("align", "center"), e.get("anchor", "middle"), wrap=True)
    return sp


def _add_table(slide, e):
    header = e.get("header", [])
    rows = e.get("rows", [])
    nrows = len(rows) + (1 if header else 0)
    ncols = e.get("ncols", len(header) or (len(rows[0]) if rows else 1))
    gs = slide.shapes.add_table(nrows, ncols, _in(e["x"]), _in(e["y"]), _in(e["w"]), _in(e["h"]))
    tbl = gs.table
    colw = e.get("col_widths")
    if colw:
        for i, w in enumerate(colw):
            if i < ncols:
                tbl.columns[i].width = _in(w)
    rowfills = e.get("rowfills", [])
    hdrfill = e.get("header_fill", "#15375F")
    def fill_cell(cell, hexcolor):
        cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(hexcolor)
    for ci, txt in enumerate(header):
        cell = tbl.cell(0, ci)
        _set_text(cell.text_frame, txt, e.get("font_size", 10.5), "#FFFFFF", True, "center", "middle")
        fill_cell(cell, hdrfill)
    for ri, row in enumerate(rows, start=1):
        for ci, txt in enumerate(row):
            cell = tbl.cell(ri, ci)
            _set_text(cell.text_frame, txt, e.get("font_size", 10.5), "#172033", False, "center", "middle")
            fill = (rowfills[ri-1][ci] if ri-1 < len(rowfills) and ci < len(rowfills[ri-1])
                    else e.get("cell_fill", "#FFFFFF"))
            fill_cell(cell, fill)
    return gs


def _add_connector(slide, e):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    _in(e["x1"]), _in(e["y1"]), _in(e["x2"]), _in(e["y2"]))
    cn.line.color.rgb = _rgb(e.get("color", "#5F6B7A"))
    cn.line.width = Pt(e.get("w", 2.0))
    ln = cn.line._get_or_add_ln()
    if e.get("arrow", True):
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    dash = e.get("dash")
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
    return cn


def _add_bar_chart(slide, e):
    cats = e.get("categories", [])
    series = e.get("series", [])
    cw, ch = _in(e["w"]), _in(e["h"])
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                    _in(e["x"]), _in(e["y"]), cw, ch)
    chart = gframe.chart
    chart.has_title = bool(e.get("title"))
    if e.get("title"):
        chart.chart_title.text_frame.text = e["title"]
    if series:
        chart.replace_data({
            "categories": cats,
            "series": series,
        })
    color = e.get("color", "#2F6FB2")
    try:
        plot = chart.plots[0]
        plot.has_data_labels = False
        for s_ in plot.series:
            s_.format.fill.solid()
            s_.format.fill.fore_color.rgb = _rgb(color)
    except Exception:
        pass
    if e.get("legend", True):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False
    return gframe


def _add_image(slide, e):
    path = e.get("path")
    if not path or not os.path.exists(path):
        raise FileNotFoundError(f"image path not found: {path}")
    slide.shapes.add_picture(path, _in(e["x"]), _in(e["y"]), _in(e["w"]), _in(e["h"]))
    return


def _add_icon(slide, e):
    """Icon placeholders are real SVG pictures added by a later COM pass (PPT can
    insert vector SVG). Here we draw a simple geometric placeholder so the layout
    is visible without the SVG pass."""
    name = e.get("name", "square")
    color = _rgb(e.get("color", "#2F6FB2"))
    x, y, w, h = _in(e["x"]), _in(e["y"]), _in(e["w"]), _in(e["h"])
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    _set_text(sp.text_frame, name[:1].upper(), min(w/Inches(1) * 18, e.get("size", 8)), "#FFFFFF",
              True, "center", "middle")
    return sp


ELEMENT_HANDLERS = {
    "text": _add_textbox,
    "box": _add_box,
    "table": _add_table,
    "connector": _add_connector,
    "bar_chart": _add_bar_chart,
    "image": _add_image,
    "icon": _add_icon,
}


def _build_slide(prs, spec, default_canvas):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # background
    bg = spec.get("background")
    if bg:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        shape.shadow.inherit = False
        _apply_fill(shape, bg)
        # send to back
        spTree = slide.shapes._spTree
        spTree.remove(shape._element)
        spTree.insert(2, shape._element)
    for e in spec.get("elements", []):
        handler = ELEMENT_HANDLERS.get(e.get("type"))
        if handler:
            try:
                handler(slide, e)
            except Exception as ex:
                print(f"  [warn] element {e.get('type')} failed: {ex}", file=sys.stderr)
    return slide


def build(scene, out_path, canvas_key="16:9"):
    w, h = CANVAS.get(canvas_key, CANVAS["16:9"])
    prs = Presentation()
    prs.slide_width = Inches(w)
    prs.slide_height = Inches(h)
    for spec in scene.get("slides", []):
        _build_slide(prs, spec, canvas_key)
    os.makedirs(os.path.dirname(os.path.abspath(out_path)) or ".", exist_ok=True)
    prs.save(out_path)
    return out_path, len(scene.get("slides", []))


def main(argv=None):
    ap = argparse.ArgumentParser(description="Build native editable .pptx from a scene JSON.")
    ap.add_argument("scene", help="scene JSON file")
    ap.add_argument("-o", "--output", default=None, help="output .pptx path")
    ap.add_argument("--canvas", default="16:9", choices=list(CANVAS.keys()))
    args = ap.parse_args(argv)
    with open(args.scene, "r", encoding="utf-8-sig") as f:
        scene = json.load(f)
    out = args.output or (os.path.splitext(args.scene)[0] + ".pptx")
    path, n = build(scene, out, args.canvas)
    print(f"OK {path}  slides={n}")


if __name__ == "__main__":
    main()
